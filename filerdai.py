import os
import cohere
import fitz 
from pptx import Presentation
from docx import Document


def convert_pptx_to_text(file):
    
    prs = Presentation(file)
    
    text = []
    
    for slide in prs.slides:
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
                
    return "\n".join(text)


def convert_docx_to_text(file):
    
    doc = Document(file)
    text = []
    
    for paragraph in doc.paragraphs:
        
        text.append(paragraph.text)
        
    return "\n".join(text)



def convert_pdf_to_text(file):
    
    doc = fitz.open(file)
    
    text = []
    
    for page in doc:
        text.append(page.get_text())
        
    return "\n".join(text)


def initialize_cohere(api_key):
    return cohere.Client(api_key)


def query_cohere(client, question, context):
    
    response = client.generate(
        model='command-xlarge-nightly',
        prompt=f"Context: {context}\n\nQuestion: {question}\n\nAnswer:",
        max_tokens=150,
        temperature=0.7,
        stop_sequences=["\n"]
    )
    
    return response.generations[0].text.strip()


def get_file_text(file_path):
    file_ext = os.path.splitext(file_path)[1].lower()

    if file_ext == ".txt":
        with open(file_path, 'r') as f:
            return f.read()
    elif file_ext == ".pptx":
        return convert_pptx_to_text(file_path)
    elif file_ext == ".docx":
        return convert_docx_to_text(file_path)
    elif file_ext == ".pdf":
        return convert_pdf_to_text(file_path)
    else:
        raise ValueError("Unsupported file format. Please use a .txt, .pptx, .docx, or .pdf file.")
    
    

def main():
    api_key = 'bP4LYKuyDiF6D4G2VQl7Bg4oYtZG2jgYpz2CtMKD'
    cohere_client = initialize_cohere(api_key)

    fname = input("Enter the file name (with extension): ")

    try:
        file_text = get_file_text(fname)
    except FileNotFoundError:
        print("File not found. Please enter a valid file name.")
        return
    except ValueError as e:
        print(e)
        return
    

    if not file_text.strip():
        print("The file is empty or could not be read.")
        return


    question = input("Enter the question you want to ask: ")


    answer = query_cohere(cohere_client, question, file_text)


    print("Answer:- ")
    print(answer)
    
    

if __name__ == "__main__":
    main()
